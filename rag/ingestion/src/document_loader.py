import os
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging

class DocumentLoader:
    """Loads documents from various file formats"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.supported_formats = {
            '.pdf': self._load_pdf,
            '.txt': self._load_txt,
            '.docx': self._load_docx,
            '.pptx': self._load_pptx,
            '.md': self._load_md
        }
    
    def load_documents(self, directory: str) -> List[Dict[str, Any]]:
        """Load all supported documents from a directory"""
        documents = []
        directory_path = Path(directory)
        
        if not directory_path.exists():
            self.logger.warning(f"Directory does not exist: {directory}")
            return documents
        
        # Find all supported files
        for file_path in directory_path.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_formats:
                try:
                    document = self.load_single_document(str(file_path))
                    if document:
                        documents.append(document)
                        self.logger.info(f"Loaded: {file_path.name}")
                except Exception as e:
                    self.logger.error(f"Failed to load {file_path}: {e}")
        
        self.logger.info(f"Loaded {len(documents)} documents from {directory}")
        return documents
    
    def load_single_document(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Load a single document"""
        file_path = Path(file_path)
        
        if not file_path.exists():
            self.logger.error(f"File does not exist: {file_path}")
            return None
        
        file_extension = file_path.suffix.lower()
        
        if file_extension not in self.supported_formats:
            self.logger.error(f"Unsupported file format: {file_extension}")
            return None
        
        try:
            # Load the document using the appropriate loader
            content = self.supported_formats[file_extension](file_path)
            
            if content:
                return {
                    'file_path': str(file_path),
                    'file_name': file_path.name,
                    'file_extension': file_extension,
                    'content': content,
                    'size': file_path.stat().st_size
                }
            else:
                self.logger.warning(f"Empty content from: {file_path}")
                return None
                
        except Exception as e:
            self.logger.error(f"Error loading {file_path}: {e}")
            return None
    
    def _load_pdf(self, file_path: Path) -> str:
        """Load PDF file"""
        try:
            import pypdf
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            self.logger.error(f"Error loading PDF {file_path}: {e}")
            return ""
    
    def _load_txt(self, file_path: Path) -> str:
        """Load text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read().strip()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read().strip()
            except Exception as e:
                self.logger.error(f"Error loading text file {file_path}: {e}")
                return ""
        except Exception as e:
            self.logger.error(f"Error loading text file {file_path}: {e}")
            return ""
    
    def _load_docx(self, file_path: Path) -> str:
        """Load DOCX file"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            self.logger.error(f"Error loading DOCX {file_path}: {e}")
            return ""
    
    def _load_pptx(self, file_path: Path) -> str:
        """Load PPTX file"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            self.logger.error(f"Error loading PPTX {file_path}: {e}")
            return ""
    
    def _load_md(self, file_path: Path) -> str:
        """Load Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read().strip()
        except Exception as e:
            self.logger.error(f"Error loading Markdown {file_path}: {e}")
            return ""
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return list(self.supported_formats.keys())
    
    def is_supported(self, file_path: str) -> bool:
        """Check if file format is supported"""
        return Path(file_path).suffix.lower() in self.supported_formats 